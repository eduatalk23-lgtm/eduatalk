"""
TimeLevelUp 프로젝트 결과보고서 PPTX 생성 스크립트
슬라이드 크기: 26.67" x 15.0" (원본 템플릿 기준, 표준 16:9의 2배)
→ 모든 좌표/크기/폰트에 2x 스케일 적용
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── 2x 스케일 헬퍼 ───
S = 2  # scale factor

def I(val):
    """Inches with 2x scale"""
    return Inches(val * S)

def P(val):
    """Pt with 2x scale"""
    return Pt(val * S)

# ─── 색상 팔레트 ───
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY = RGBColor(0x0F, 0x1A, 0x33)
BLUE = RGBColor(0x2E, 0x86, 0xDE)
LIGHT_BLUE = RGBColor(0x5B, 0xB5, 0xF5)
ACCENT_BLUE = RGBColor(0x3A, 0x7B, 0xD5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xEA, 0xED)
MID_GRAY = RGBColor(0x9A, 0x9A, 0x9A)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
RED = RGBColor(0xEF, 0x44, 0x44)

# ─── 슬라이드 크기 (원본 템플릿) ───
SLIDE_WIDTH = Emu(24384000)   # 26.67"
SLIDE_HEIGHT = Emu(13716000)  # 15.0"


def add_bg(slide, color=DARK_NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def rrect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height, text, size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """텍스트 박스 — size는 논리 pt (자동 2x 스케일)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = P(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "맑은 고딕"
    p.alignment = align
    return txBox


def add_p(tf, text, size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = P(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "맑은 고딕"
    p.alignment = align
    p.space_before = P(space)
    return p


def header(slide, num, title):
    """공통 섹션 헤더"""
    rect(slide, I(0), I(0), SLIDE_WIDTH, I(0.04), BLUE)
    tb(slide, I(0.4), I(0.15), I(8), I(0.35),
       f"{num}. {title}", size=17, color=LIGHT_BLUE, bold=True)


def cnum(slide, left, top, num, sz=0.35, color=BLUE, fsz=11):
    """번호 원형"""
    c = circle(slide, I(left), I(top), I(sz), color)
    c.text_frame.paragraphs[0].text = str(num)
    c.text_frame.paragraphs[0].font.size = P(fsz)
    c.text_frame.paragraphs[0].font.color.rgb = WHITE
    c.text_frame.paragraphs[0].font.bold = True
    c.text_frame.paragraphs[0].font.name = "맑은 고딕"
    c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return c


# ─── 프레젠테이션 생성 ───
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank = prs.slide_layouts[6]


# ════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
rect(s, I(0), I(0), SLIDE_WIDTH, I(0.04), BLUE)
rect(s, I(0), I(7.46), SLIDE_WIDTH, I(0.04), BLUE)
rect(s, I(0.5), I(0.6), I(0.05), I(4.8), BLUE)

tb(s, I(0.75), I(0.6), I(8), I(0.4),
   "AI리더 1기 프로젝트 결과보고서", size=12, color=LIGHT_BLUE)

tb(s, I(0.75), I(1.1), I(10), I(0.7),
   "TimeLevelUp", size=36, color=WHITE, bold=True)

tb(s, I(0.75), I(1.9), I(10), I(0.5),
   "AI 기반 맞춤형 학습 관리 플랫폼", size=18, color=LIGHT_BLUE)

rect(s, I(0.75), I(2.5), I(2.5), I(0.03), BLUE)

tb(s, I(0.75), I(2.7), I(8), I(0.35),
   "팀명 : TimeLevelUp", size=12, color=WHITE)
tb(s, I(0.75), I(3.1), I(8), I(0.35),
   "우조현", size=12, color=MID_GRAY)

tags = ["Next.js 16", "TypeScript 5", "Supabase", "Gemini AI", "Toss Payments"]
for i, tag in enumerate(tags):
    x = 7.5 + (i % 3) * 1.8
    y = 5.2 + (i // 3) * 0.45
    shape = rrect(s, I(x), I(y), I(1.6), I(0.35), NAVY)
    shape.text_frame.paragraphs[0].text = tag
    shape.text_frame.paragraphs[0].font.size = P(9)
    shape.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    shape.text_frame.paragraphs[0].font.name = "맑은 고딕"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════
# 슬라이드 2: 목차
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
rect(s, I(0), I(0), SLIDE_WIDTH, I(0.04), BLUE)

tb(s, I(0.4), I(0.2), I(3), I(0.35),
   "OUTLINE", size=10, color=MID_GRAY, bold=True)

items = [
    ("01", "프로젝트 개요", "OUTLINE"),
    ("02", "프로젝트 팀 구성 및 역할", "TEAM & ROLE"),
    ("03", "프로젝트 수행 절차 및 방법", "PROCEDURE"),
    ("04", "프로젝트 수행 경과", "PROGRESS"),
    ("05", "자체 평가 의견", "ASSESSMENT"),
]

for i, (num, title, eng) in enumerate(items):
    y = 0.7 + i * 1.2

    cnum(s, 0.5, y + 0.02, num, 0.35, BLUE, 11)

    tb(s, I(1.1), I(y), I(7), I(0.35),
       f"{num}. {title}", size=14, color=WHITE, bold=True)

    tb(s, I(9.5), I(y + 0.05), I(3), I(0.25),
       eng, size=8, color=MID_GRAY)

    if i < len(items) - 1:
        rect(s, I(1.1), I(y + 0.55), I(11.5), I(0.01),
             RGBColor(0x2A, 0x3A, 0x5A))


# ════════════════════════════════════════
# 슬라이드 3: 01. 프로젝트 개요
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "01", "프로젝트 개요")

# 좌측: 주제 및 배경
tb(s, I(0.4), I(0.6), I(5.8), I(0.35),
   "프로젝트 주제 및 선정 배경", size=14, color=WHITE, bold=True)

txBox = tb(s, I(0.4), I(1.0), I(5.8), I(2.5), "", size=10, color=LIGHT_GRAY)
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "한국 학원(사교육) 시장에서 학생 개인별 맞춤 학습 관리가 부재한 문제를 해결하기 위한 AI 기반 학습 관리 플랫폼"
p.font.size = P(10)
p.font.color.rgb = LIGHT_GRAY
p.font.name = "맑은 고딕"

add_p(tf, "", size=4, color=LIGHT_GRAY, space=4)
add_p(tf, "• 기존 학원 시스템: 단순 출결/결제 중심, 학습 분석 부재", size=10, color=LIGHT_GRAY, space=4)
add_p(tf, "• 내신/모의고사 분석 → AI 맞춤 학습 계획 자동 생성", size=10, color=LIGHT_GRAY, space=4)
add_p(tf, "• 학생·학부모·관리자 3자 연결 통합 플랫폼", size=10, color=LIGHT_GRAY, space=4)
add_p(tf, "• 멀티 테넌트 아키텍처로 프랜차이즈 학원 지원", size=10, color=LIGHT_GRAY, space=4)

# 좌측 하단: 기대 효과
tb(s, I(0.4), I(3.8), I(5.8), I(0.35),
   "활용방안 및 기대 효과", size=13, color=BLUE, bold=True)

effects = [
    "• 학원 운영 효율 극대화 — 출결/일정/결제 자동화",
    "• AI 성적 기반 최적 학습 계획으로 학생 맞춤 학습",
    "• 학부모 실시간 현황 공유로 신뢰 구축",
    "• CRM 리드 관리로 학원 매출 성장 지원",
]
for i, eff in enumerate(effects):
    tb(s, I(0.4), I(4.25) + I(i * 0.35), I(5.8), I(0.35),
       eff, size=10, color=LIGHT_GRAY)

# 우측: 활용 기술
card = rrect(s, I(6.6), I(0.6), I(6.0), I(5.6), NAVY)
rect(s, I(6.6), I(0.6), I(6.0), I(0.04), BLUE)

tb(s, I(6.85), I(0.72), I(5.5), I(0.35),
   "활용 장비 및 기술", size=13, color=BLUE, bold=True)

tech = [
    ("Frontend", "Next.js 16, React 19, TypeScript 5, Tailwind CSS 4"),
    ("Backend", "Supabase (PostgreSQL + RLS + Realtime)"),
    ("AI / ML", "Gemini AI, Claude, OpenAI (멀티 프로바이더)"),
    ("결제", "Toss Payments (토스페이먼츠 SDK)"),
    ("메시징", "Ppurio SMS / 알림톡, Resend 이메일"),
    ("기타", "Google Calendar API, QR 출결, Vercel 배포"),
]
for i, (cat, detail) in enumerate(tech):
    y = 1.2 + i * 0.45
    tb(s, I(6.85), I(y), I(1.3), I(0.35),
       f"▸ {cat}", size=9, color=LIGHT_BLUE, bold=True)
    tb(s, I(8.2), I(y), I(4.2), I(0.35),
       detail, size=9, color=LIGHT_GRAY)

# 우측 하단: 프로젝트 내용
tb(s, I(6.85), I(4.0), I(5.5), I(0.35),
   "프로젝트 내용", size=13, color=BLUE, bold=True)

contents = [
    "• 4개 역할 시스템: 학생, 관리자, 학부모, 슈퍼관리자",
    "• AI 학습 계획 생성 (Gemini/Claude/OpenAI)",
    "• 콜드 스타트 콘텐츠 추천 (신규 학생 대응)",
    "• 성적 분석 백분위 엔진 (내신/모의고사)",
    "• CRM, SMS, 토스 결제, 구글 캘린더 연동",
]
for i, item in enumerate(contents):
    tb(s, I(6.85), I(4.45) + I(i * 0.35), I(5.5), I(0.35),
       item, size=9, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 4: 02. 팀 구성 및 역할
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "02", "프로젝트 팀 구성 및 역할")

tb(s, I(0.4), I(0.55), I(11), I(0.3),
   "해당 프로젝트를 진행하면서 주도적으로 참여한 부분을 중심으로 작성합니다.",
   size=10, color=MID_GRAY)

rrect(s, I(0.3), I(0.95), I(12.7), I(6.3), NAVY)

# 프로필
c = circle(s, I(0.6), I(1.1), I(0.7), BLUE)
c.text_frame.paragraphs[0].text = "우"
c.text_frame.paragraphs[0].font.size = P(22)
c.text_frame.paragraphs[0].font.color.rgb = WHITE
c.text_frame.paragraphs[0].font.bold = True
c.text_frame.paragraphs[0].font.name = "맑은 고딕"
c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

tb(s, I(1.5), I(1.1), I(5), I(0.35),
   "우조현", size=18, color=WHITE, bold=True)
tb(s, I(1.5), I(1.5), I(6), I(0.3),
   "Full-Stack Developer & Project Lead", size=11, color=LIGHT_BLUE)

# 4개 역할 카드
roles = [
    ("기획 & 디자인", BLUE, [
        "서비스 기획 및 UX 설계",
        "학원 도메인 분석",
        "요구사항 정의",
    ]),
    ("프론트엔드", ACCENT_BLUE, [
        "Next.js 16 App Router",
        "50+ 페이지 (4개 역할)",
        "DnD 학습 플래너",
    ]),
    ("백엔드 & DB", PURPLE, [
        "Supabase 71개 마이그레이션",
        "RLS 보안 정책 설계",
        "멀티 테넌트 아키텍처",
    ]),
    ("AI & 외부 연동", GREEN, [
        "AI 멀티 프로바이더 파이프라인",
        "토스 결제 / 뿌리오 SMS",
        "Google Calendar / CI/CD",
    ]),
]

for i, (title, color, items_list) in enumerate(roles):
    x = 0.5 + i * 3.15
    y = 2.2

    card = rrect(s, I(x), I(y), I(2.95), I(4.8), DARK_NAVY)
    rect(s, I(x), I(y), I(2.95), I(0.04), color)

    tb(s, I(x + 0.15), I(y + 0.15), I(2.65), I(0.35),
       title, size=12, color=color, bold=True)

    for j, item in enumerate(items_list):
        tb(s, I(x + 0.15), I(y + 0.65 + j * 0.45), I(2.65), I(0.4),
           f"• {item}", size=10, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 5: 03. 수행 절차 및 방법
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "03", "프로젝트 수행 절차 및 방법")

tb(s, I(0.4), I(0.55), I(7), I(0.3),
   "프로젝트 사전 기획부터 MVP 완성까지의 수행 절차입니다.",
   size=10, color=MID_GRAY)

# 총 개발 기간
period = rrect(s, I(8.8), I(0.18), I(3.8), I(0.5), BLUE)
tb(s, I(8.9), I(0.2), I(3.6), I(0.2),
   "총 개발 기간", size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(s, I(8.9), I(0.4), I(3.6), I(0.2),
   "2025. 11. 20 ~ 2026. 02. 19", size=9, color=WHITE, align=PP_ALIGN.CENTER)

# 연결선
rect(s, I(0.8), I(1.65), I(11.5), I(0.02), RGBColor(0x2A, 0x3A, 0x5A))

phases = [
    ("사전 기획", "11/20 ~ 12/05", BLUE,
     ["도메인 분석 및 요구사항 정의", "서비스 기획안 작성", "DB 스키마 설계"]),
    ("핵심 기능 구현", "12/05 ~ 01/10", ACCENT_BLUE,
     ["인증/인가 (역할 기반)", "학습 플래너 (위자드+DnD)", "성적 입력/분석 엔진"]),
    ("AI 시스템 구축", "01/10 ~ 01/25", PURPLE,
     ["콜드 스타트 파이프라인", "AI 학습 계획 생성", "배치 처리 자동화"]),
    ("외부 연동", "01/25 ~ 02/10", ORANGE,
     ["토스 결제 시스템", "뿌리오 SMS/알림톡", "Google Calendar 동기화"]),
    ("MVP 완성", "02/10 ~ 02/19", GREEN,
     ["CRM 영업 관리", "학부모 포털", "성능 최적화 및 배포"]),
]

for i, (title, period_txt, color, items_list) in enumerate(phases):
    x = 0.3 + i * 2.5

    cnum(s, x + 0.7, 1.48, str(i + 1), 0.35, color, 11)

    tb(s, I(x), I(1.95), I(2.3), I(0.4),
       title, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    tb(s, I(x), I(2.35), I(2.3), I(0.25),
       period_txt, size=8, color=color, align=PP_ALIGN.CENTER)

    card = rrect(s, I(x), I(2.7), I(2.3), I(4.5), NAVY)
    for j, item in enumerate(items_list):
        tb(s, I(x + 0.12), I(2.85 + j * 0.5), I(2.1), I(0.45),
           f"• {item}", size=9, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 6: 04. 시스템 아키텍처
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "04", "프로젝트 수행 경과")
tb(s, I(0.4), I(0.55), I(10), I(0.35),
   "시스템 아키텍처", size=16, color=WHITE, bold=True)

layers = [
    ("클라이언트 (Frontend)", BLUE,
     "Next.js 16 App Router  |  React 19  |  TypeScript 5  |  Tailwind CSS 4  |  React Query  |  DnD Kit  |  Recharts"),
    ("서버 (Backend)", ACCENT_BLUE,
     "Server Actions  |  API Routes  |  역할 기반 Auth  |  도메인: Plan, Score, SMS, Payment, CRM, Chat, Attendance"),
    ("AI / ML 레이어", PURPLE,
     "Gemini AI  |  Claude (Anthropic)  |  OpenAI  |  콜드 스타트 파이프라인  |  학습 계획 생성  |  백분위 엔진"),
    ("데이터베이스 & 외부 서비스", GREEN,
     "Supabase PostgreSQL (RLS)  |  Realtime  |  Toss Payments  |  Ppurio SMS/알림톡  |  Google Calendar"),
    ("인프라 & CI/CD", ORANGE,
     "Vercel 배포  |  GitHub Actions  |  pnpm  |  ESLint  |  Vitest  |  자동 배포 (main push)  |  일일 배치"),
]

for i, (name, color, detail) in enumerate(layers):
    y = 1.05 + i * 1.2

    rrect(s, I(0.3), I(y), I(12.7), I(1.1), NAVY)
    rect(s, I(0.3), I(y), I(0.06), I(1.1), color)

    tb(s, I(0.55), I(y + 0.08), I(5), I(0.35),
       name, size=12, color=color, bold=True)

    tb(s, I(0.55), I(y + 0.5), I(12.2), I(0.4),
       detail, size=10, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 7: 04. AI 학습 시스템
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "04", "프로젝트 수행 경과")
tb(s, I(0.4), I(0.55), I(11), I(0.35),
   "핵심 기능 ① — AI 학습 계획 & 콜드 스타트 시스템", size=14, color=WHITE, bold=True)

tb(s, I(0.4), I(1.0), I(8), I(0.3),
   "콜드 스타트 파이프라인 (Cold Start Pipeline)", size=12, color=BLUE, bold=True)

pipeline = [
    ("입력 검증", "교과/과목/난이도\n콘텐츠 타입", BLUE),
    ("쿼리 생성", "Gemini AI로\n검색 쿼리 생성", ACCENT_BLUE),
    ("콘텐츠 검색", "웹 검색 +\nDB 캐시 조회", PURPLE),
    ("결과 파싱", "AI 응답 파싱\n및 구조화", ORANGE),
    ("추천 완료", "랭킹/필터링\nDB 저장 반환", GREEN),
]

for i, (step, desc, color) in enumerate(pipeline):
    x = 0.25 + i * 2.5

    card = rrect(s, I(x), I(1.4), I(2.2), I(1.5), NAVY)
    rect(s, I(x), I(1.4), I(2.2), I(0.04), color)

    tb(s, I(x + 0.1), I(1.5), I(2.0), I(0.3),
       step, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    tb(s, I(x + 0.1), I(1.85), I(2.0), I(0.8),
       desc, size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    if i < len(pipeline) - 1:
        tb(s, I(x + 2.15), I(1.7), I(0.4), I(0.4),
           "→", size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 멀티 프로바이더
tb(s, I(0.4), I(3.15), I(8), I(0.3),
   "멀티 프로바이더 AI 시스템", size=12, color=BLUE, bold=True)

providers = [
    ("Gemini AI", "콘텐츠 검색/추천, 콜드 스타트", "Free Tier: 일 20회/분 15회 자동 관리", BLUE),
    ("Claude (Anthropic)", "학습 계획 생성, 하이브리드 플랜", "스트리밍 지원, 부분 재생성 가능", PURPLE),
    ("OpenAI", "대체 프로바이더, 플랜 최적화", "Fallback 용도, 유연한 확장성", GREEN),
]

for i, (name, feat, note, color) in enumerate(providers):
    x = 0.25 + i * 4.2

    card = rrect(s, I(x), I(3.55), I(3.9), I(2.0), NAVY)
    rect(s, I(x), I(3.55), I(0.05), I(2.0), color)

    tb(s, I(x + 0.2), I(3.65), I(3.5), I(0.3),
       name, size=11, color=color, bold=True)
    tb(s, I(x + 0.2), I(4.05), I(3.5), I(0.4),
       feat, size=9, color=LIGHT_GRAY)
    tb(s, I(x + 0.2), I(4.5), I(3.5), I(0.4),
       note, size=8, color=MID_GRAY)

# 하단: Rate Limit & AI 계획 생성
tb(s, I(0.4), I(5.8), I(5), I(0.3),
   "Rate Limit 보호 & 배치 처리", size=11, color=BLUE, bold=True)

rl_card = rrect(s, I(0.25), I(6.15), I(6.0), I(1.2), NAVY)
rl_items = [
    "• Gemini Free Tier 할당량 자동 추적 (일 20회, 분 15회)",
    "• API 한도 초과 시 DB 캐시 Fallback 자동 전환",
    "• GitHub Actions 배치: 매일 새벽 3시 (KST)",
]
for i, item in enumerate(rl_items):
    tb(s, I(0.4), I(6.2 + i * 0.35), I(5.6), I(0.3),
       item, size=9, color=LIGHT_GRAY)

tb(s, I(6.6), I(5.8), I(5), I(0.3),
   "AI 학습 계획 생성 기능", size=11, color=BLUE, bold=True)

plan_card = rrect(s, I(6.5), I(6.15), I(6.2), I(1.2), NAVY)
plan_items = [
    "• 위자드 기반 다단계 계획 생성 (과목/기간/교재)",
    "• 하이브리드 플랜: AI 추천 + 사용자 커스터마이징",
    "• 스트리밍 생성 / 부분 재생성 지원",
]
for i, item in enumerate(plan_items):
    tb(s, I(6.7), I(6.2 + i * 0.35), I(5.8), I(0.3),
       item, size=9, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 8: 04. 성적 분석 / 결제
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "04", "프로젝트 수행 경과")
tb(s, I(0.4), I(0.55), I(11), I(0.35),
   "핵심 기능 ② — 성적 분석 엔진 / 토스 결제 시스템", size=14, color=WHITE, bold=True)

# 좌측: 성적
card_l = rrect(s, I(0.25), I(1.0), I(6.2), I(6.2), NAVY)
rect(s, I(0.25), I(1.0), I(6.2), I(0.04), BLUE)

tb(s, I(0.5), I(1.15), I(5.8), I(0.35),
   "성적 분석 엔진", size=14, color=BLUE, bold=True)

score_items = [
    "• 내신 / 모의고사 성적 입력 및 관리",
    "• 백분위 자동 계산 엔진",
    "• Recharts 기반 성적 추이 시각화 차트",
    "• 등급 변환 규칙 엔진",
    "• 학교별 / 전국 벤치마크 비교",
    "• 과목별 분석 및 취약점 파악",
]
for i, item in enumerate(score_items):
    tb(s, I(0.5), I(1.65 + i * 0.55), I(5.8), I(0.45),
       item, size=11, color=LIGHT_GRAY)

# 우측: 결제
card_r = rrect(s, I(6.8), I(1.0), I(6.2), I(6.2), NAVY)
rect(s, I(6.8), I(1.0), I(6.2), I(0.04), ORANGE)

tb(s, I(7.05), I(1.15), I(5.8), I(0.35),
   "토스 결제 시스템", size=14, color=ORANGE, bold=True)

pay_items = [
    "• Toss Payments SDK 연동",
    "• 일회 결제 + 배치 결제 지원",
    "• Webhook 기반 결제 확인 (안전 검증)",
    "• 멱등성 키로 중복 결제 방지",
    "• 할인 시스템 (프로그램별 할인율)",
    "• 월별 매출 리포트 (프로그램별 집계)",
]
for i, item in enumerate(pay_items):
    tb(s, I(7.05), I(1.65 + i * 0.55), I(5.8), I(0.45),
       item, size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 9: 04. SMS / CRM
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "04", "프로젝트 수행 경과")
tb(s, I(0.4), I(0.55), I(11), I(0.35),
   "핵심 기능 ③ — SMS 메시징 / CRM 영업 관리", size=14, color=WHITE, bold=True)

# 좌측: SMS
card_l = rrect(s, I(0.25), I(1.0), I(6.2), I(6.2), NAVY)
rect(s, I(0.25), I(1.0), I(6.2), I(0.04), GREEN)

tb(s, I(0.5), I(1.15), I(5.8), I(0.35),
   "SMS & 알림톡", size=14, color=GREEN, bold=True)

sms_items = [
    "• Ppurio 뿌리오 API 연동 (SMS/LMS)",
    "• 카카오 알림톡 비즈 메시지 발송",
    "• 예약 발송 및 취소 기능",
    "• 발송 결과 동기화 및 추적",
    "• 상담 D-1 자동 리마인더",
    "• 출결/결제/상담 트랜잭션 문자",
]
for i, item in enumerate(sms_items):
    tb(s, I(0.5), I(1.65 + i * 0.55), I(5.8), I(0.45),
       item, size=11, color=LIGHT_GRAY)

# 우측: CRM
card_r = rrect(s, I(6.8), I(1.0), I(6.2), I(6.2), NAVY)
rect(s, I(6.8), I(1.0), I(6.2), I(0.04), PURPLE)

tb(s, I(7.05), I(1.15), I(5.8), I(0.35),
   "CRM 영업 관리", size=14, color=PURPLE, bold=True)

crm_items = [
    "• 세일즈 리드 파이프라인 관리",
    "• 리드 스코어링 시스템 (자동 점수화)",
    "• 활동 로그 및 태스크 관리",
    "• 상담사 배정 시스템",
    "• 프로그램 연결 (PRO/PRE/MAS 등)",
    "• 알림톡 연동 자동 알림 발송",
]
for i, item in enumerate(crm_items):
    tb(s, I(7.05), I(1.65 + i * 0.55), I(5.8), I(0.45),
       item, size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 10: 04. 역할별 기능 & 성과
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "04", "프로젝트 수행 경과")
tb(s, I(0.4), I(0.55), I(11), I(0.35),
   "역할별 사용자 기능 & 프로젝트 규모", size=14, color=WHITE, bold=True)

role_data = [
    ("학생", BLUE, ["대시보드, 오늘의 학습(3-Dock)", "성적 입력/분석, 학습 패턴", "캠프, 채팅, 습관, 출결"]),
    ("관리자", ORANGE, ["학생/성적/상담 관리", "SMS 센터, 출결 관리", "CRM 영업, 결제/수납"]),
    ("학부모", GREEN, ["자녀 대시보드, 성적 조회", "결제/수납, 학습 리포트", "학습 이력, 목표 설정"]),
    ("슈퍼관리자", PURPLE, ["멀티 테넌트(지점) 관리", "어드민 사용자 관리", "교육과정, 약관 관리"]),
]

for i, (role, color, items_list) in enumerate(role_data):
    x = 0.25 + i * 3.2
    card = rrect(s, I(x), I(1.0), I(3.0), I(3.0), NAVY)
    rect(s, I(x), I(1.0), I(3.0), I(0.04), color)

    tb(s, I(x + 0.15), I(1.12), I(2.7), I(0.35),
       role, size=12, color=color, bold=True)

    for j, item in enumerate(items_list):
        tb(s, I(x + 0.15), I(1.6 + j * 0.45), I(2.7), I(0.4),
           f"• {item}", size=9, color=LIGHT_GRAY)

# 성과 수치
tb(s, I(0.4), I(4.2), I(4), I(0.3),
   "프로젝트 규모", size=12, color=BLUE, bold=True)

metrics = [
    ("50+", "페이지", BLUE),
    ("71", "DB 마이그레이션", ACCENT_BLUE),
    ("458+", "Git 커밋", PURPLE),
    ("3", "AI 프로바이더", ORANGE),
    ("10+", "외부 서비스 연동", GREEN),
]

for i, (num, label, color) in enumerate(metrics):
    x = 0.25 + i * 2.5
    card = rrect(s, I(x), I(4.6), I(2.3), I(2.6), NAVY)

    tb(s, I(x), I(4.75), I(2.3), I(0.6),
       num, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    tb(s, I(x), I(5.5), I(2.3), I(0.35),
       label, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 슬라이드 11: 04. 시연 영상
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "04", "프로젝트 수행 경과")

txBox = tb(s, I(0.4), I(0.55), I(12), I(0.6), "", size=10, color=MID_GRAY)
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "세부 기능 소개, 화면 구동 및 기능 동작 여부를 시연영상으로 제작합니다."
p.font.size = P(10)
p.font.color.rgb = MID_GRAY
p.font.name = "맑은 고딕"
add_p(tf, "- 용량 제한 : 5 ~ 10분 내(100MB 이하), 기능별 소개 음성 포함", size=9, color=MID_GRAY)

demo_card = rrect(s, I(1.0), I(1.5), I(11.3), I(5.7), NAVY)

tb(s, I(2.0), I(1.8), I(9), I(0.5),
   "시연 영상 촬영 예정 항목", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

demo_items = [
    ("1", "학생 로그인 → 대시보드 → 오늘의 학습 (3-Dock 플래너)", BLUE),
    ("2", "AI 학습 계획 생성 (위자드 → AI 추천 → 결과 확인)", ACCENT_BLUE),
    ("3", "성적 입력 및 분석 차트 (내신/모의고사 백분위)", PURPLE),
    ("4", "관리자 CRM 영업 관리 & SMS 발송", ORANGE),
    ("5", "학부모 포털 (성적 조회 / 결제)", GREEN),
    ("6", "QR 출결 체크인 & 상담 일정 관리", RED),
]

for i, (num, desc, color) in enumerate(demo_items):
    y = 2.5 + i * 0.7

    cnum(s, 2.0, y, num, 0.35, color, 11)

    tb(s, I(2.6), I(y - 0.02), I(9), I(0.45),
       desc, size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 12: 05. 자체 평가 의견
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
header(s, "05", "자체 평가 의견")
tb(s, I(0.4), I(0.55), I(11), I(0.3),
   "프로젝트 결과물에 대한 완성도, 달성도, 느낀 점을 평가합니다.",
   size=10, color=MID_GRAY)

# 좌상: 완성도
card = rrect(s, I(0.25), I(1.0), I(6.2), I(2.8), NAVY)
rect(s, I(0.25), I(1.0), I(6.2), I(0.04), BLUE)
tb(s, I(0.5), I(1.12), I(4), I(0.35),
   "완성도 평가 (10점 만점)", size=12, color=BLUE, bold=True)
tb(s, I(4.5), I(1.1), I(1.7), I(0.4),
   "8 / 10", size=18, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)

eval_items = [
    "핵심 기능(학습 계획, 성적 분석, 출결, 결제, SMS) 모두 구현",
    "AI 멀티 프로바이더, 멀티 테넌트 등 고급 아키텍처 적용",
    "일부 UI 폴리싱 및 E2E 테스트 커버리지 보완 필요",
]
for i, item in enumerate(eval_items):
    tb(s, I(0.5), I(1.6 + i * 0.4), I(5.7), I(0.35),
       f"• {item}", size=9, color=LIGHT_GRAY)

# 우상: 잘한 부분
card = rrect(s, I(6.8), I(1.0), I(6.2), I(2.8), NAVY)
rect(s, I(6.8), I(1.0), I(6.2), I(0.04), GREEN)
tb(s, I(7.05), I(1.12), I(5.5), I(0.35),
   "잘한 부분", size=12, color=GREEN, bold=True)

good_items = [
    "학원 도메인에 특화된 실무 지향 설계",
    "AI 콜드 스타트 시스템으로 신규 학생 즉시 대응",
    "71개 마이그레이션 / RLS 기반 보안 아키텍처",
    "실 결제/SMS 연동으로 실무 수준 완성도",
]
for i, item in enumerate(good_items):
    tb(s, I(7.05), I(1.55 + i * 0.4), I(5.7), I(0.35),
       f"• {item}", size=9, color=LIGHT_GRAY)

# 좌하: 느낀 점
card = rrect(s, I(0.25), I(4.0), I(6.2), I(3.2), NAVY)
rect(s, I(0.25), I(4.0), I(6.2), I(0.04), PURPLE)
tb(s, I(0.5), I(4.12), I(5.5), I(0.35),
   "느낀 점 & 성과", size=12, color=PURPLE, bold=True)

feel_items = [
    "1인 풀스택으로 기획~배포 전 과정 수행, 실무 역량 성장",
    "AI Rate Limit 관리와 비용 최적화 전략 학습",
    "Supabase RLS 멀티 테넌트 아키텍처 설계 경험",
    "외부 서비스 연동 시 보안과 안정성의 중요성 체감",
    "실 사용자 서비스 구축 경험 (학원 현장 피드백 반영)",
]
for i, item in enumerate(feel_items):
    tb(s, I(0.5), I(4.6 + i * 0.4), I(5.7), I(0.35),
       f"• {item}", size=9, color=LIGHT_GRAY)

# 우하: 개선점
card = rrect(s, I(6.8), I(4.0), I(6.2), I(3.2), NAVY)
rect(s, I(6.8), I(4.0), I(6.2), I(0.04), ORANGE)
tb(s, I(7.05), I(4.12), I(5.5), I(0.35),
   "향후 개선점", size=12, color=ORANGE, bold=True)

improve_items = [
    "E2E 테스트 (Playwright) 커버리지 확대",
    "성능 모니터링 및 에러 트래킹 도입 (Sentry 등)",
    "AI 학습 계획 정확도 향상 (사용자 피드백 루프)",
    "PWA 오프라인 지원 강화",
    "프랜차이즈 실 론칭을 위한 부하 테스트",
]
for i, item in enumerate(improve_items):
    tb(s, I(7.05), I(4.6 + i * 0.4), I(5.7), I(0.35),
       f"• {item}", size=9, color=LIGHT_GRAY)


# ════════════════════════════════════════
# 슬라이드 13: 감사합니다
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
rect(s, I(0), I(0), SLIDE_WIDTH, I(0.04), BLUE)
rect(s, I(0), I(7.46), SLIDE_WIDTH, I(0.04), BLUE)

tb(s, I(0), I(2.5), SLIDE_WIDTH, I(0.8),
   "감사합니다", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

tb(s, I(0), I(3.4), SLIDE_WIDTH, I(0.5),
   "Thank You", size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

rect(s, I(5.0), I(4.1), I(3.3), I(0.03), BLUE)

tb(s, I(0), I(4.3), SLIDE_WIDTH, I(0.4),
   "TimeLevelUp — AI 기반 맞춤형 학습 관리 플랫폼", size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

tb(s, I(0), I(4.8), SLIDE_WIDTH, I(0.4),
   "우조현  |  AI리더 1기", size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ─── 저장 ───
output_path = "/Users/johyeon-u/Desktop/결과보고서_TimeLevelUp(AI기반_맞춤형_학습관리_플랫폼).pptx"
prs.save(output_path)
print(f"✅ 결과보고서 저장 완료: {output_path}")
print(f"📊 총 슬라이드 수: {len(prs.slides)}")
print(f"📐 슬라이드 크기: {prs.slide_width / 914400:.1f}\" x {prs.slide_height / 914400:.1f}\"")
